"""Enhance the ECE 143 presentation with EDA figures, XGBoost results, and pipeline details."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path
import copy

SRC = Path("/Users/msnarainshriraam/Downloads/ECE 143 Presentation.pptx")
OUT = Path("/Users/msnarainshriraam/Downloads/ECE 143 Presentation Enhanced.pptx")
FIG = Path("/Users/msnarainshriraam/Documents/ECE_143_Project/ece143-final-project-1/churn-prevention-agents/figures")

prs = Presentation(str(SRC))

DARK = RGBColor(0x2D, 0x2D, 0x2D)
ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
BLUE = RGBColor(0x34, 0x98, 0xDB)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

W = prs.slide_width
H = prs.slide_height


def _add_textbox(slide, left, top, width, height, text, font_name="Lato",
                 font_size=Pt(14), bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def _add_bullet_slide(slide, left, top, width, height, bullets, font_name="Lato",
                      font_size=Pt(13), color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.space_after = Pt(6)
        p.level = 0
        for run in p.runs:
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = color


def _blank_slide(prs):
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    for sl in prs.slide_layouts:
        if sl.name in ("Blank", "BLANK"):
            layout = sl
            break
    return prs.slides.add_slide(layout)


def _get_layout(prs, name):
    for sl in prs.slide_layouts:
        if name.lower() in sl.name.lower():
            return sl
    return prs.slide_layouts[0]


NS = '{http://schemas.openxmlformats.org/presentationml/2006/main}'

def _move_slide_to(prs, slide, target_idx):
    """Move the last slide in sldIdLst to target_idx."""
    sldIdLst = prs.element.find(f'{NS}sldIdLst')
    items = list(sldIdLst)
    item = items[-1]
    sldIdLst.remove(item)
    remaining = list(sldIdLst)
    if target_idx >= len(remaining):
        sldIdLst.append(item)
    else:
        remaining[target_idx].addprevious(item)


# ════════════════════════════════════════════
#  NEW SLIDES TO INSERT
# ════════════════════════════════════════════

insert_after = 5  # after Slide 5 (Data Overview), 0-indexed

# ── SLIDE A: EDA — Churn Distribution + Membership ──
slideA = _blank_slide(prs)
_add_textbox(slideA, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
             "EDA: What Drives Churn?", "Montserrat", Pt(28), True, DARK, PP_ALIGN.CENTER)

slideA.shapes.add_picture(str(FIG / "01_churn_distribution.png"),
                          Inches(0.3), Inches(1.0), Inches(4.2), Inches(2.8))
slideA.shapes.add_picture(str(FIG / "02_churn_by_membership.png"),
                          Inches(4.8), Inches(1.0), Inches(5.0), Inches(2.8))

_add_bullet_slide(slideA, Inches(0.3), Inches(3.9), Inches(9.2), Inches(1.5), [
    "• 54.1% overall churn rate (20,012 churned vs 16,980 retained) — balanced for ML",
    "• Membership category is the #1 predictor: categories 0 & 2 → ~97% churn; 3 & 4 → 0% churn",
    "• Upgrading customers to higher membership tiers is the highest-ROI retention lever",
], "Lato", Pt(11), DARK)

_move_slide_to(prs, slideA, insert_after + 1)

# ── SLIDE B: EDA — Numeric Features + Correlation ──
slideB = _blank_slide(prs)
_add_textbox(slideB, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
             "EDA: Feature Analysis", "Montserrat", Pt(28), True, DARK, PP_ALIGN.CENTER)

slideB.shapes.add_picture(str(FIG / "03_numeric_by_churn.png"),
                          Inches(0.2), Inches(0.9), Inches(5.5), Inches(2.6))
slideB.shapes.add_picture(str(FIG / "04_correlation_heatmap.png"),
                          Inches(5.5), Inches(0.9), Inches(4.3), Inches(3.5))

_add_bullet_slide(slideB, Inches(0.2), Inches(3.7), Inches(9.5), Inches(1.8), [
    "• Retained customers have higher transaction values (median $30.6K vs $25.4K) and more wallet points (749 vs 648)",
    "• Age, tenure, days since login show virtually no difference — NOT useful churn predictors",
    "• support_risk and past_complaint are perfectly correlated (r=1.0) — one is redundant",
    "• Monetary engagement separates churners from retained, not demographics or recency",
], "Lato", Pt(11), DARK)

_move_slide_to(prs, slideB, insert_after + 2)

# ── SLIDE C: EDA — Segments + Feedback ──
slideC = _blank_slide(prs)
_add_textbox(slideC, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
             "EDA: Segments & Feedback Signals", "Montserrat", Pt(28), True, DARK, PP_ALIGN.CENTER)

slideC.shapes.add_picture(str(FIG / "05_churn_by_segments.png"),
                          Inches(0.2), Inches(0.9), Inches(5.3), Inches(2.2))
slideC.shapes.add_picture(str(FIG / "06_complaint_support_churn.png"),
                          Inches(5.3), Inches(0.9), Inches(4.5), Inches(2.0))

_add_bullet_slide(slideC, Inches(0.2), Inches(3.2), Inches(9.5), Inches(2.2), [
    "• Value segmentation matters: low/medium value → 58-59% churn; high value → 45% (below average)",
    "• Engagement segmentation barely differentiates churn (all ~53-55%) — needs rework",
    "• Feedback is a near-perfect signal: types 4-6, 8 → 0% churn; types 0-3, 7 → 63-65% churn",
    "• Past complaints have surprisingly small effect on churn (+1.7% only)",
    "• Key insight: act immediately on negative feedback — it's the clearest early warning",
], "Lato", Pt(11), DARK)

_move_slide_to(prs, slideC, insert_after + 3)

# ── SLIDE D: EDA — Tenure + Price Sensitivity ──
slideD = _blank_slide(prs)
_add_textbox(slideD, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
             "EDA: Tenure & Price Sensitivity", "Montserrat", Pt(28), True, DARK, PP_ALIGN.CENTER)

slideD.shapes.add_picture(str(FIG / "07_tenure_distribution.png"),
                          Inches(0.3), Inches(0.9), Inches(5.0), Inches(2.8))
slideD.shapes.add_picture(str(FIG / "08_churn_by_price_sensitivity.png"),
                          Inches(5.3), Inches(0.9), Inches(4.3), Inches(2.8))

_add_bullet_slide(slideD, Inches(0.3), Inches(3.9), Inches(9.2), Inches(1.5), [
    "• Churn is spread uniformly across all tenure levels — not a lifecycle problem",
    "• Unlike many subscription businesses, long-tenured customers are NOT meaningfully safer",
    "• Price sensitivity: <2% difference — not a meaningful churn differentiator in this dataset",
    "• Churn drivers are behavioral/engagement-based, not demographic or price-based",
], "Lato", Pt(11), DARK)

_move_slide_to(prs, slideD, insert_after + 4)

# ── SLIDE E: XGBoost Training Results ──
slideE = _blank_slide(prs)
_add_textbox(slideE, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
             "XGBoost Model: Training Results", "Montserrat", Pt(28), True, DARK, PP_ALIGN.CENTER)

_add_textbox(slideE, Inches(0.5), Inches(0.85), Inches(4), Inches(0.4),
             "Model Performance", "Montserrat", Pt(18), True, BLUE, PP_ALIGN.LEFT)

metrics_table = slideE.shapes.add_table(6, 2, Inches(0.5), Inches(1.3), Inches(3.5), Inches(2.5)).table
metrics_table.columns[0].width = Inches(1.8)
metrics_table.columns[1].width = Inches(1.7)
metrics_data = [("Metric", "Value"), ("Accuracy", "93.35%"), ("Precision", "93.14%"),
                ("Recall", "94.68%"), ("F1-Score", "93.90%"), ("AUC", "97.59%")]
for row_idx, (k, v) in enumerate(metrics_data):
    for col_idx, val in enumerate([k, v]):
        cell = metrics_table.cell(row_idx, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Lato"
                r.font.size = Pt(13)
                r.font.bold = row_idx == 0 or (row_idx == 5 and col_idx == 1)
                r.font.color.rgb = WHITE if row_idx == 0 else (ACCENT if row_idx == 5 and col_idx == 1 else DARK)
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x34, 0x49, 0x5E)

_add_textbox(slideE, Inches(5.0), Inches(0.85), Inches(4.5), Inches(0.4),
             "Feature Importances (Top 5)", "Montserrat", Pt(18), True, BLUE, PP_ALIGN.LEFT)

feat_table = slideE.shapes.add_table(6, 2, Inches(5.0), Inches(1.3), Inches(4.5), Inches(2.5)).table
feat_table.columns[0].width = Inches(2.8)
feat_table.columns[1].width = Inches(1.7)
feat_data = [("Feature", "Importance"), ("membership_category", "46.4%"),
             ("points_in_wallet", "21.0%"), ("feedback", "4.3%"),
             ("avg_transaction_value", "2.5%"), ("All others", "< 1.5% each")]
for row_idx, (k, v) in enumerate(feat_data):
    for col_idx, val in enumerate([k, v]):
        cell = feat_table.cell(row_idx, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Consolas" if col_idx == 0 and row_idx > 0 else "Lato"
                r.font.size = Pt(13)
                r.font.bold = row_idx == 0 or (row_idx in (1, 2) and col_idx == 1)
                r.font.color.rgb = WHITE if row_idx == 0 else (ACCENT if row_idx in (1, 2) and col_idx == 1 else DARK)
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x34, 0x49, 0x5E)

_add_bullet_slide(slideE, Inches(0.5), Inches(4.0), Inches(9.0), Inches(1.3), [
    "• XGBClassifier: 300 trees, max_depth=6, learning_rate=0.1, class-balanced via scale_pos_weight",
    "• EDA confirmed: membership_category (46%) + points_in_wallet (21%) = 67% of model decisions",
    "• Trained via /api/v1/train endpoint, logged to MLflow for experiment tracking",
], "Lato", Pt(11), DARK)

# Insert after the existing model training slides (slide 9 original = index 8, but we added 4 so it's 12)
_move_slide_to(prs, slideE, insert_after + 5 + 4)  # after slide 9 shifted

# ── SLIDE F: Agentic Pipeline Debate ──
slideF = _blank_slide(prs)
_add_textbox(slideF, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
             "Agentic Pipeline: The Debate Mechanism", "Montserrat", Pt(28), True, DARK, PP_ALIGN.CENTER)

_add_textbox(slideF, Inches(0.3), Inches(0.9), Inches(4.3), Inches(0.35),
             "Round 1: Strategist Proposes", "Montserrat", Pt(16), True, BLUE)
_add_bullet_slide(slideF, Inches(0.3), Inches(1.3), Inches(4.3), Inches(1.3), [
    "• Generic email campaigns & loyalty discounts",
    "• Personalized points showcasing",
    "• FAQ guides for rewards programs",
], "Lato", Pt(11), DARK)

_add_textbox(slideF, Inches(0.3), Inches(2.5), Inches(4.3), Inches(0.35),
             "Critic Rejects → Rating: 3/10", "Montserrat", Pt(16), True, ACCENT)
_add_bullet_slide(slideF, Inches(0.3), Inches(2.9), Inches(4.3), Inches(1.3), [
    "✗ No measurement strategy or KPIs defined",
    "✗ No control groups for A/B testing",
    "✗ Root cause 'avg' too generic for targeted action",
], "Lato", Pt(11), ACCENT)

_add_textbox(slideF, Inches(5.2), Inches(0.9), Inches(4.5), Inches(0.35),
             "Round 2: Strategist Revises", "Montserrat", Pt(16), True, BLUE)
_add_bullet_slide(slideF, Inches(5.2), Inches(1.3), Inches(4.5), Inches(1.3), [
    "• Validation-first: survey to check model calibration",
    "• A/B test: 10% discount vs feature highlight",
    "• Personalized points redemption with control group",
], "Lato", Pt(11), DARK)

_add_textbox(slideF, Inches(5.2), Inches(2.5), Inches(4.5), Inches(0.35),
             "Critic Approves → Rating: 8/10", "Montserrat", Pt(16), True, GREEN)
_add_bullet_slide(slideF, Inches(5.2), Inches(2.9), Inches(4.5), Inches(1.3), [
    "✓ Measurement strategy with specific KPIs",
    "✓ A/B testing with control groups included",
    "✓ Phased rollout to minimize risk",
], "Lato", Pt(11), GREEN)

_add_textbox(slideF, Inches(0.3), Inches(4.2), Inches(9.2), Inches(0.8),
             "The system self-improves through adversarial debate — the Critic forces the Strategist to produce rigorous, measurable plans.",
             "Lato", Pt(13), True, DARK, PP_ALIGN.CENTER)

# Place near end, before thank you
_move_slide_to(prs, slideF, len(list(prs.element.find(f'{NS}sldIdLst'))) - 2)

# ── SLIDE G: Pipeline Results ──
slideG = _blank_slide(prs)
_add_textbox(slideG, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
             "Pipeline Output: Approved Retention Plan", "Montserrat", Pt(28), True, DARK, PP_ALIGN.CENTER)

_add_textbox(slideG, Inches(0.3), Inches(0.85), Inches(9.2), Inches(0.3),
             "19,941 at-risk customers identified by XGBoost → segmented → strategies generated and debated",
             "Lato", Pt(13), False, GRAY, PP_ALIGN.CENTER)

_add_textbox(slideG, Inches(0.3), Inches(1.3), Inches(4.5), Inches(0.35),
             "Segment 1: Transaction Patterns (18,056 customers)", "Montserrat", Pt(14), True, ACCENT)
_add_bullet_slide(slideG, Inches(0.3), Inches(1.7), Inches(4.5), Inches(1.6), [
    "Priority: HIGH | Avg CLV: $16,637",
    "1. Validate 95% churn probability via customer survey",
    "2. A/B test: 10% discount vs personalized feature highlight",
    "3. Analyze qualitative feedback to decompose root cause",
], "Lato", Pt(11), DARK)

_add_textbox(slideG, Inches(5.2), Inches(1.3), Inches(4.5), Inches(0.35),
             "Segment 2: Loyalty Points (1,885 customers)", "Montserrat", Pt(14), True, BLUE)
_add_bullet_slide(slideG, Inches(5.2), Inches(1.7), Inches(4.5), Inches(1.6), [
    "Priority: MEDIUM | Avg CLV: $1,972",
    "1. Validate churn probability via targeted survey",
    "2. A/B test: personalized redemption vs standard reminder",
    "3. Phased rollout starting with small subset",
], "Lato", Pt(11), DARK)

_add_textbox(slideG, Inches(0.3), Inches(3.4), Inches(9.2), Inches(0.35),
             "Risks Flagged by the Agents", "Montserrat", Pt(14), True, RGBColor(0xE6, 0x7E, 0x22))
_add_bullet_slide(slideG, Inches(0.3), Inches(3.8), Inches(9.2), Inches(1.4), [
    "• 95%+ churn probabilities suggest model calibration review needed before full deployment",
    "• No historical intervention data — all strategies are experimental, requiring careful monitoring",
    "• 'avg' root cause needs decomposition; A/B tests need minimum sample sizes for significance",
], "Lato", Pt(11), DARK)

_move_slide_to(prs, slideG, len(list(prs.element.find(f'{NS}sldIdLst'))) - 2)

# ── SLIDE H: Key Takeaways ──
slideH = _blank_slide(prs)
_add_textbox(slideH, Inches(0.5), Inches(0.2), Inches(9), Inches(0.6),
             "Key Takeaways", "Montserrat", Pt(28), True, DARK, PP_ALIGN.CENTER)

takeaways = [
    ("1", "Membership tier is the #1 churn lever", "46% of model importance — upgrade programs have highest ROI", ACCENT),
    ("2", "XGBoost achieves 97.6% AUC", "Confirms EDA: membership + wallet points = 67% of decisions", BLUE),
    ("3", "Feedback is a near-perfect early warning", "Types 4-6, 8 → 0% churn; types 0-3, 7 → 63% churn", GREEN),
    ("4", "Agentic pipeline bridges prediction → action", "Multi-agent debate ensures strategy quality (3/10 → 8/10)", RGBColor(0x9B, 0x59, 0xB6)),
    ("5", "End-to-end: data → model → agents → campaigns", "Fully automated, API-driven, Docker-deployed, MLflow-tracked", RGBColor(0xE6, 0x7E, 0x22)),
]

for i, (num, title, desc, color) in enumerate(takeaways):
    y = Inches(0.95 + i * 0.85)
    _add_textbox(slideH, Inches(0.5), y, Inches(0.5), Inches(0.4),
                 num, "Montserrat", Pt(22), True, color, PP_ALIGN.CENTER)
    _add_textbox(slideH, Inches(1.1), y, Inches(8.3), Inches(0.35),
                 title, "Montserrat", Pt(16), True, DARK)
    _add_textbox(slideH, Inches(1.1), Emu(y + Inches(0.35)), Inches(8.3), Inches(0.35),
                 desc, "Lato", Pt(12), False, GRAY)

_move_slide_to(prs, slideH, len(list(prs.element.find(f'{NS}sldIdLst'))) - 2)

# ═══════════════════════════════════
#  Fill empty Slide 16 & 17
# ═══════════════════════════════════
# Slide 16 (VIP At Risk vs Standard) and 17 are mostly empty
# They are now shifted due to insertions. Let's find them by checking for empty content.
# We'll leave them as-is since the new slides cover the content better.

prs.save(str(OUT))
print(f"Enhanced presentation saved to: {OUT}")
print(f"Total slides: {len(prs.slides)}")
